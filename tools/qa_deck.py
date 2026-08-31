# -*- coding: utf-8 -*-
"""기획서 레이아웃 자동 검수 — 슬라이드 밖 이탈 / 텍스트 넘침 / 겹침 탐지"""
import sys, math, unicodedata
from pptx import Presentation
from pptx.util import Emu

EMU = 914400.0
SW, SH = 13.333, 7.5
TOL = 0.02

def wide(ch):
    return unicodedata.east_asian_width(ch) in ("W", "F")

def text_width_pt(s, size):
    """한글=size, 영문/숫자≈0.52*size 로 근사한 문자열 폭(pt)"""
    return sum(size if wide(c) else size * 0.52 for c in s)

def needed_height_in(tf, box_w_in):
    """텍스트프레임이 실제로 필요로 하는 높이(inch) 추정"""
    total_pt, box_w_pt = 0.0, box_w_in * 72
    if box_w_pt <= 0:
        return 0
    for p in tf.paragraphs:
        s = "".join(r.text for r in p.runs)
        size = max((r.font.size.pt for r in p.runs if r.font.size), default=12)
        ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.2
        sb = p.space_before.pt if p.space_before else 0
        lines = max(1, math.ceil(text_width_pt(s, size) / box_w_pt)) if s else 1
        total_pt += lines * size * ls * 1.06 + sb
    return total_pt / 72

def rects_overlap(a, b, pad=0.0):
    ax, ay, aw, ah = a; bx, by, bw, bh = b
    return not (ax + aw <= bx + pad or bx + bw <= ax + pad or
                ay + ah <= by + pad or by + bh <= ay + pad)

prs = Presentation(sys.argv[1])
issues = []
for i, sl in enumerate(prs.slides, 1):
    texts = []
    for sh in sl.shapes:
        if sh.left is None or sh.top is None:
            continue
        x, y = sh.left / EMU, sh.top / EMU
        w, h = (sh.width or 0) / EMU, (sh.height or 0) / EMU
        name = sh.shape_type
        # 1) 슬라이드 밖 이탈
        if x < -TOL or y < -TOL or x + w > SW + TOL or y + h > SH + TOL:
            issues.append((i, "OUT", f"{name} at ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f} "
                                     f"→ 우/하단 {x+w:.2f},{y+h:.2f}"))
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        need = needed_height_in(sh.text_frame, w)
        snippet = sh.text_frame.text.strip().replace("\n", " / ")[:42]
        # 2) 텍스트 넘침 (박스 높이 대비)
        if need > h * 1.25 + 0.06:
            issues.append((i, "OVERFLOW",
                           f"h={h:.2f} 필요≈{need:.2f} | {snippet}"))
        # 3) 슬라이드 하단 이탈 (텍스트가 실제로 차지하는 높이 기준)
        if y + need > SH + 0.04:
            issues.append((i, "BOTTOM",
                           f"y={y:.2f}+need{need:.2f}={y+need:.2f} > {SH} | {snippet}"))
        texts.append((x, y, w, max(h, need), snippet))
    # 4) 텍스트끼리 겹침
    for a in range(len(texts)):
        for b in range(a + 1, len(texts)):
            ra, rb = texts[a][:4], texts[b][:4]
            if rects_overlap(ra, rb, pad=0.03):
                ov_h = min(ra[1]+ra[3], rb[1]+rb[3]) - max(ra[1], rb[1])
                if ov_h > 0.10:
                    issues.append((i, "COLLIDE",
                                   f"{ov_h:.2f}in | '{texts[a][4]}' ↔ '{texts[b][4]}'"))

print(f"슬라이드 {len(prs.slides)}장 · 발견 {len(issues)}건\n")
from collections import Counter
print(Counter(k for _, k, _ in issues), "\n")
cur = None
for s, k, m in issues:
    if s != cur:
        print(f"\n── S{s}"); cur = s
    print(f"  [{k}] {m}")
