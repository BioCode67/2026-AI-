# -*- coding: utf-8 -*-
"""기획서 PPTX 자동 생성 — outputs/ 의 실제 산출물 수치를 그대로 반영"""
from __future__ import annotations
import json, sys
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
sys.path.insert(0, str(Path(__file__).resolve().parent))
from config import TAB, FIG, DELIV

FONT = "맑은 고딕"
INK, MUTE, ACC = C(0x15, 0x1C, 0x26), C(0x5D, 0x68, 0x74), C(0xD6, 0x45, 0x45)
BLUE, GREEN, AMBER = C(0x2E, 0x7D, 0xD1), C(0x3E, 0x9B, 0x7C), C(0xC9, 0x8A, 0x1E)
LIGHT, LINE, WHITE = C(0xF5, 0xF7, 0xFA), C(0xE0, 0xE5, 0xEB), C(0xFF, 0xFF, 0xFF)
W, H, M = 13.333, 7.5, 0.72


def tb(sl, x, y, w, h, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
       line=1.35, space=0, font=FONT):
    s = sl.shapes.add_textbox(In(x), In(y), In(w), In(h))
    f = s.text_frame; f.word_wrap = True; f.margin_left = f.margin_right = 0
    f.margin_top = f.margin_bottom = 0
    for i, ln in enumerate(str(text).split("\n")):
        p = f.paragraphs[0] if i == 0 else f.add_paragraph()
        p.alignment = align; p.line_spacing = line
        if i: p.space_before = Pt(space)
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        r.font.name = font
    return s


def rect(sl, x, y, w, h, fill=LIGHT, lc=None, radius=True):
    sh = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        In(x), In(y), In(w), In(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if lc: sh.line.color.rgb = lc; sh.line.width = Pt(1)
    else:  sh.line.fill.background()
    sh.shadow.inherit = False
    if radius:
        try: sh.adjustments[0] = 0.06
        except Exception: pass
    sh.text_frame.text = ""
    return sh


def slide(prs, eyebrow=None, title=None, lede=None):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    y = 0.52
    if eyebrow:
        tb(sl, M, y, W - 2 * M, .3, eyebrow, 11.5, True, ACC); y += .36
    if title:
        tb(sl, M, y, W - 2 * M, .62, title, 27, True, INK, line=1.18); y += .78
    if lede:
        tb(sl, M, y, W - 2 * M, .45, lede, 13.5, False, MUTE, line=1.45)
    return sl


def pic(sl, name, x, y, w=None, h=None):
    p = FIG / name
    if not p.exists(): return None
    return sl.shapes.add_picture(str(p), In(x), In(y),
                                 In(w) if w else None, In(h) if h else None)


def fit(sl, name, x, y, bw, bh):
    """박스 (bw×bh) 안에 비율 유지하며 가운데 배치"""
    from PIL import Image
    p = FIG / name
    if not p.exists(): return
    try:
        iw, ih = Image.open(p).size; ar = iw / ih
    except Exception:
        ar = 1.5
    w, h = bw, bw / ar
    if h > bh: h, w = bh, bh * ar
    sl.shapes.add_picture(str(p), In(x + (bw - w) / 2), In(y + (bh - h) / 2), In(w), In(h))


def bullets(sl, x, y, w, items, size=13.5, gap=.42, marker="▪"):
    """항목 높이를 누적해 배치한다(줄 수가 다른 본문이 겹치지 않도록)."""
    yy = y
    for head, body in items:
        tb(sl, x, yy, .25, .3, marker, size, True, ACC)
        tb(sl, x + .3, yy, w - .3, .28, head, size, True, INK)
        h = .3
        if body:
            bs = size - 1.5
            lines = sum(max(1, int(len(ln) * bs / ((w - .3) * 72) * 1.02) + 1)
                        for ln in body.split("\n"))
            bh = lines * bs * 1.42 * 1.06 / 72
            tb(sl, x + .3, yy + .3, w - .3, bh, body, bs, False, MUTE, line=1.42)
            h += bh
        yy += h + gap * .38


def kpi_card(sl, x, y, w, h, lab, val, sub, accent=ACC, vsize=24):
    h = max(h, 1.34)                       # 라벨+값(24pt)+부제 2줄이 들어갈 최소 높이
    rect(sl, x, y, w, h, WHITE, LINE)
    rect(sl, x, y, .075, h, accent, radius=False)
    tb(sl, x + .28, y + .16, w - .5, .22, lab, 10.5, True, MUTE)
    tb(sl, x + .28, y + .42, w - .5, .40, val, vsize, True, INK)
    tb(sl, x + .28, y + .92, w - .5, .38, sub, 10, False, MUTE, line=1.32)


def tablette(sl, df, x, y, w, colw=None, size=10.5, maxrows=10, head_bg=LIGHT,
             rh=.335):
    df = df.head(maxrows)
    nr, nc = len(df) + 1, len(df.columns)
    t = sl.shapes.add_table(nr, nc, In(x), In(y), In(w), In(rh * nr)).table
    if colw:
        tot = sum(colw)
        for i, c in enumerate(colw):
            t.columns[i].width = Emu(int(In(w) * c / tot))
    for j, cname in enumerate(df.columns):
        cell = t.cell(0, j); cell.text = str(cname)
        cell.fill.solid(); cell.fill.fore_color.rgb = head_bg
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]; r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = INK; r.font.name = FONT
    for i, row in enumerate(df.itertuples(index=False), start=1):
        for j, v in enumerate(row):
            cell = t.cell(i, j); cell.text = str(v)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
            r = p.runs[0]; r.font.size = Pt(size - .5)
            r.font.color.rgb = INK; r.font.name = FONT
    for i in range(nr):
        t.rows[i].height = In(rh)
    return t


def footer(sl, n, txt="천안 균형발전 나침반(CBC)"):
    tb(sl, M, H - .46, 6, .3, txt, 9.5, False, MUTE)
    tb(sl, W - M - 1, H - .46, 1, .3, str(n), 9.5, False, MUTE, align=PP_ALIGN.RIGHT)


def build(team="○○팀", members="홍길동"):
    S = json.loads((TAB / "09_요약지표.json").read_text(encoding="utf-8"))
    prov = pd.read_csv(TAB / "00_데이터_출처.csv")
    cbi  = pd.read_csv(TAB / "02_CBI_균형발전지수.csv")
    _ewp = TAB / "04_쇠퇴조기경보_예측.csv"
    ew   = pd.read_csv(_ewp) if _ewp.exists() else pd.DataFrame(
        columns=["zone", "risk", "risk_pred", "risk_delta"])
    sites = pd.read_csv(TAB / "08_생활SOC_투자우선순위.csv")
    real = S["전체실데이터"]

    prs = Presentation()
    prs.slide_width, prs.slide_height = In(W), In(H)
    n = 0

    # ══ 1. 표지 ══════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, C(0x11, 0x18, 0x22), radius=False)
    rect(sl, 0, 0, .13, H, ACC, radius=False)
    tb(sl, 1.15, 1.55, 11, .35, "2026년 천안시 AI·데이터 기반 정책 아이디어 경진대회",
       13, True, C(0xE0, 0x8A, 0x8A))
    tb(sl, 1.15, 2.08, 11.2, .95, "천안 균형발전 나침반", 54, True, WHITE, line=1.1)
    tb(sl, 1.15, 3.12, 11.2, .5, "Cheonan Balance Compass", 20, False,
       C(0x8A, 0x97, 0xA8))
    rect(sl, 1.15, 3.85, 5.6, .045, ACC, radius=False)
    tb(sl, 1.15, 4.12, 11, 1.0,
       "천안시민 누구나 걸어서 닿는 생활환경을 바라며,\n공공데이터로 먼저 살펴본 균형발전 이야기",
       23, True, C(0xE9, 0xED, 0xF2), line=1.35)
    for i, (lab, val) in enumerate([("지정과제", "③ 지역균형발전"),
                                    ("응모분야", "AI 모델 개발"),
                                    ("분석단위", f"{S['생활권수']}개 생활권 · 지표 {S['지표수']}종")]):
        x = 1.15 + i * 3.55
        tb(sl, x, 5.72, 3.3, .25, lab, 10.5, True, C(0x6E, 0x7B, 0x8C))
        tb(sl, x, 6.0, 3.3, .3, val, 14, True, WHITE)
    tb(sl, 1.15, 6.72, 5.4, .3, f"{team} · {members}", 12, False, C(0x8A, 0x97, 0xA8))
    if not real:
        tb(sl, W - M - 5.2, 6.72, 5.2, .3,
           "※ 일부 지표 예시 데이터 — 실데이터 재실행 후 제출", 10, True, AMBER,
           align=PP_ALIGN.RIGHT)

    # ══ 2. 한 장 요약 ════════════════════════════════════════
    n += 1
    sl = slide(prs, "EXECUTIVE SUMMARY", "한 장으로 보는 제안",
               "천안시민의 생활여건 차이를 데이터로 살펴보고, 도움이 조금 더 필요한 곳을 "
               "먼저 찾아보고자 준비했습니다.")
    cards = [("격차 배율", f'{S["격차배율"]}배', f'신도심 {S["신도심_CBI"]} vs 원도심 {S["원도심_CBI"]}', ACC),
             ("불균등도(지니)", f'{S["지니계수"]}', f'{S["생활권수"]}개 생활권 CBI 기준', AMBER),
             ("조기경보 방향 적중률",
              (f'{S["방향적중률"]:.0f}%' if S.get("방향적중률") == S.get("방향적중률")
               else "—") if S.get("예측엔진", True) else "데이터 대기",
              f'베이스라인 대비 MAE {S["MAE_개선율"]:+.1f}%' if S.get("예측엔진", True)
              else "인구 파일 투입 시 활성화", BLUE),
             ("투자 수혜", f'{S["신규수혜인구"]:,}명', f'{S["추천입지수"]}개소 · 커버리지 +{S["커버리지개선"]}%p', GREEN)]
    for i, (a, b, c, col) in enumerate(cards):
        kpi_card(sl, M + i * 3.06, 2.26, 2.86, 1.40, a, b, c, col, vsize=22)
    cols = [("저희가 주목한 것", "한 도시 안에 생긴 생활여건의 차이",
             "서북부 신도심이 커지는 동안 동남부 원도심과 농촌면은\n"
             "반대 방향으로 움직였습니다. 같은 천안시민인데 사시는 곳에\n"
             "따라 일상의 여건이 달라지고 있었습니다."),
            ("저희가 준비한 것", "살펴보기 → 미리 알기 → 순서 정하기",
             "공공데이터 12개 지표로 생활권별 여건을 수치화하고,\n"
             "어느 동네가 어려워질지 미리 살펴본 뒤,\n"
             "생활SOC를 어디부터 놓으면 좋을지 순서를 제안드립니다."),
            ("이렇게 쓰이면 좋겠습니다", "담당자분께 드리는 근거 자료",
             f"생활권별 지수와 유형, 3년 뒤 전망과 그 이유, "
             f"{S['추천입지수']}개소 우선순위표입니다.\n"
             f"스크립트 한 줄이면 매년 새 데이터로 다시 만드실 수 있습니다.")]
    for i, (tag, head, body) in enumerate(cols):
        x = M + i * 4.09
        rect(sl, x, 4.02, 3.89, 2.62, WHITE, LINE)
        rect(sl, x, 4.02, 3.89, .075, [ACC, BLUE, GREEN][i], radius=False)
        tb(sl, x + .28, 4.28, 3.3, .25, tag, 11, True, [ACC, BLUE, GREEN][i])
        tb(sl, x + .28, 4.58, 3.35, .5, head, 15.5, True, INK, line=1.25)
        tb(sl, x + .28, 5.24, 3.35, 1.3, body, 11.5, False, MUTE, line=1.5)
    footer(sl, n)

    # ══ 3. 문제 정의 ═════════════════════════════════════════
    n += 1
    sl = slide(prs, "BACKGROUND", "한 도시 안에 서로 다른 시간이 흐르고 있습니다",
               "수도권 전철·KTX와 산업단지를 축으로 서북부가 커지는 사이, 동남부 원도심과 농촌면은 "
               "반대 방향으로 움직였습니다.")
    lo = cbi.nsmallest(5, "CBI")[["zone", "CBI"]]
    hi = cbi.nlargest(3, "CBI")[["zone", "CBI"]]
    rect(sl, M, 2.42, 6.0, 3.9, WHITE, LINE)
    tb(sl, M + .32, 2.66, 5.4, .3, "공공데이터로 확인한 차이", 14.5, True, INK)
    facts = [(f'{S["격차배율"]}배', f'신도심 평균 CBI {S["신도심_CBI"]} vs 원도심 평균 {S["원도심_CBI"]}'),
             (f'{S["지니계수"]}', f'25개 생활권 CBI 지니계수 — 변동계수 {S["변동계수"]}%'),
             (f'{S["최고"].split()[0]} ↔ {S["최저"].split()[0]}',
              f'최상위 {S["최고"].split()[1]}점 ↔ 최하위 {S["최저"].split()[1]}점, 같은 시(市)'),
             (f'{len(cbi[cbi.stage == "쇠퇴"])}개 생활권', '여러 지표가 함께 낮은 유형으로 분류되었습니다')]
    for i, (big, small) in enumerate(facts):
        y = 3.12 + i * .78
        tb(sl, M + .32, y, 2.55, .38, big, 19, True, ACC)
        tb(sl, M + .32, y + .38, 5.4, .3, small, 11.5, False, MUTE, line=1.35)
    rect(sl, M + 6.35, 2.42, 5.55, 3.9, LIGHT)
    tb(sl, M + 6.65, 2.66, 5, .3, "왜 지금 살펴보면 좋을까요", 14.5, True, INK)
    bullets(sl, M + 6.65, 3.14, 4.95, [
        ("변화는 서서히, 그러나 되돌리기는 어렵게 옵니다",
         "상권이 비면 오가는 분이 줄고, 그래서 더 비는 흐름이 생깁니다.\n"
         "많이 진행된 뒤에 쓰는 예산은 같은 금액으로도 효과가 적습니다."),
        ("조짐은 눈에 잘 띄지 않습니다",
         "통계가 뚜렷하게 나빠졌을 때는 이미 상당히 진행된 뒤인 경우가 많습니다.\n"
         "미리 알아차릴 신호가 하나쯤 더 있으면 좋겠다는 생각에서 출발했습니다."),
        ("천안시는 이런 관찰에 좋은 조건을 갖췄습니다",
         "신도심·기성시가지·원도심·읍·농촌면이 한 시(市) 안에 모두 있어,\n"
         "동네 유형별로 무엇이 필요한지 함께 살펴볼 수 있습니다."),
    ], size=12.5, gap=.4)
    footer(sl, n)

    # ══ 4. 기존 접근의 한계 ═══════════════════════════════════
    n += 1
    sl = slide(prs, "OUR ROLE", "저희가 조심스럽게 보태고 싶은 부분",
               "이미 잘 해오고 계신 일들 위에, 공공데이터가 거들 수 있는 자리를 찾아보았습니다.")
    comp = pd.DataFrame({
        "업무 장면": ["사업지 검토", "살펴보는 시점", "예산 배분", "선정 사유 설명"],
        "이미 하고 계신 일 (가장 중요한 부분)":
            ["현장을 다니며 쌓은 경험과 주민 의견 청취",
             "민원과 현장 변화를 통한 상황 파악",
             "지역 여건과 형평성을 고려한 배분",
             "사업 필요성에 대한 행정적 판단"],
        "데이터가 거들 수 있는 일":
            [f"{S['지표수']}개 지표로 계산한 참고용 수치 한 장",
             "지표 변화로 미리 살펴보는 3년 뒤 전망",
             "한 곳 더 지을 때의 효과를 순서대로 계산",
             "판단의 근거 수치를 표로 함께 제시"]})
    tablette(sl, comp, M, 2.62, W - 2 * M, colw=[1.2, 3.0, 2.8], size=11.5, maxrows=4,
             rh=.42)
    rect(sl, M, 4.86, W - 2 * M, 1.42, LIGHT)
    tb(sl, M + .35, 5.08, 11.4, .32,
       "데이터는 현장의 경험을 대신할 수 없습니다. 다만 한 번 더 확인해 보는 참고자료는 될 수 있다고 생각합니다.",
       14.5, True, INK)
    tb(sl, M + .35, 5.5, 11.4, .68,
       "새로운 데이터를 더 사자는 제안이 아닙니다. 천안시가 이미 쓰실 수 있는 공공데이터만 엮어 세 가지 참고자료를 만들었고,\n"
       "추가 비용이나 시스템 없이 스크립트 한 줄로 매년 갱신하실 수 있게 준비했습니다.",
       12.5, False, MUTE, line=1.5)
    footer(sl, n)

    # ══ 5. 아키텍처 ══════════════════════════════════════════
    n += 1
    sl = slide(prs, "ARCHITECTURE", "3-엔진 파이프라인",
               "공공데이터 입력부터 정책 산출물까지 단일 스크립트로 연결된다.")
    stages = [("INPUT", "공공데이터", "주민등록 인구\nLOCALDATA 인허가\n상가정보·생활SOC\n빈집·버스정류장", MUTE),
              ("① DIAGNOSE", "진단 엔진", "엔트로피 가중 CBI\nK-means 유형화\n5대 도메인 분해", ACC),
              ("② PREDICT", "예측 엔진", "LightGBM 회귀\nLeave-One-Zone-Out CV\nSHAP 요인분해", BLUE),
              ("③ PRESCRIBE", "처방 엔진", "MCLP 탐욕 최적화\n형평성 제약\n한계효용 곡선", GREEN),
              ("OUTPUT", "정책 산출물", "생활권 CBI 지도\n3년 후 위험 경보\n투자 우선순위표", INK)]
    bw, gap = 2.24, .28
    for i, (tag, head, body, col) in enumerate(stages):
        x = M + i * (bw + gap)
        rect(sl, x, 2.66, bw, 2.55, WHITE, LINE)
        rect(sl, x, 2.66, bw, .075, col, radius=False)
        tb(sl, x + .2, 2.88, bw - .35, .25, tag, 9.5, True, col)
        tb(sl, x + .2, 3.18, bw - .35, .32, head, 14.5, True, INK)
        tb(sl, x + .2, 3.62, bw - .35, 1.4, body, 11, False, MUTE, line=1.55)
        if i < 4:
            tb(sl, x + bw + .02, 3.72, .26, .3, "▶", 13, True, C(0xB8, 0xC0, 0xCA))
    rect(sl, M, 5.52, W - 2 * M, .95, LIGHT)
    tb(sl, M + .35, 5.74, 11.4, .55,
       "설계 원칙 ① 재현성 — 모든 단계가 공개 스크립트(`python3 src/run_all.py`)로 동일 재현\n"
       "설계 원칙 ② 정직성 — 실데이터/예시데이터 상태를 로그·산출물에 항상 표기해 오인 제출을 차단",
       12.5, False, INK, line=1.55)
    footer(sl, n)

    # ══ 6. 활용 데이터 명세 ═══════════════════════════════════
    n += 1
    sl = slide(prs, "DATA", "활용 데이터 명세",
               "전부 정식 공공 포털에서 합법적으로 내려받은 파일이며, 개인식별정보를 포함하지 않는다.")
    ds = pd.DataFrame({
        "데이터셋": ["주민등록 인구통계(연령별)", "지방행정 인허가데이터", "소상공인 상가(상권)정보",
                  "생활SOC 표준데이터 6종", "빈집 통계", "버스정류소 현황"],
        "제공기관": ["행정안전부", "행정안전부", "소상공인시장진흥공단",
                  "공공데이터포털", "통계청 / 한국부동산원", "천안시"],
        "출처 URL": ["jumin.mois.go.kr", "localdata.go.kr", "data.go.kr",
                   "data.go.kr", "kosis.kr / emptyhomes.kr", "cheonan.go.kr"],
        "주요 컬럼": ["행정구역, 총인구수, 연령대별 인구", "인허가일자, 폐업일자, 영업상태명, 소재지주소, 업태",
                   "상권업종분류, 행정동명, 위경도", "시설명, 소재지주소, 위경도",
                   "시군구, 빈집수, 빈집사유", "정류소명, 소재지"],
        "산출 지표": ["인구증감률·청년비율·고령비율", "상권 순증감·폐업률·업종다양성",
                   "사업체밀도·생활편의시설", "생활SOC 접근성·인구당 SOC수",
                   "빈집률", "정류장 밀도"]})
    tablette(sl, ds, M, 2.52, W - 2 * M, colw=[2.0, 1.3, 1.5, 2.6, 2.2], size=10.5, maxrows=6)
    tb(sl, M, 4.86, W - 2 * M, .3, "수집 시점 및 데이터 상태", 13.5, True, INK)
    pv = prov.copy()
    pv["상태"] = pv["상태"].map({"REAL": "실데이터", "ILLUSTRATIVE": "예시(미투입)"})
    tablette(sl, pv[["지표군", "상태", "출처", "비고"]], M, 5.2, W - 2 * M,
             colw=[1.1, 1.3, 3.4, 3.4], size=10, maxrows=6)
    footer(sl, n)

    # ══ 7. 공간단위 설계 ═════════════════════════════════════
    n += 1
    sl = slide(prs, "METHOD · 01", "분석 공간단위를 25개 '생활권'으로 재설계한 이유",
               "데이터마다 공간 키가 달라서 생기는 오차를, 지표를 만들기 전에 먼저 제거했다.")
    rect(sl, M, 2.55, 5.75, 1.72, C(0xFD, 0xF2, 0xF2))
    tb(sl, M + .3, 2.78, 5.2, .3, "문제 — 공간 키 불일치", 13.5, True, ACC)
    tb(sl, M + .3, 3.16, 5.2, 1.0,
       "· 주민등록 인구통계는 행정동 기준 (성정1동 / 성정2동)\n"
       "· 인허가·상가정보는 주소 문자열 = 법정동 기준 (성정동 하나)\n"
       "· 주소만으로는 성정동 점포를 1동·2동으로 나눌 근거가 없다",
       11.5, False, INK, line=1.6)
    rect(sl, M + 6.1, 2.55, 5.8, 1.72, C(0xF0, 0xF7, 0xF2))
    tb(sl, M + 6.4, 2.78, 5.2, .3, "해결 — 1:1 대응이 되는 수준까지 통합", 13.5, True, GREEN)
    tb(sl, M + 6.4, 3.16, 5.2, 1.0,
       "· 성정1·2동 → 성정동 / 쌍용1·2·3동 → 쌍용동\n"
       "· 원성1·2동 → 원성동 / 불당1·2동 → 불당동 / 부성1·2동 → 부성동\n"
       "· 결과: 31개 행정동 → 분석 가능한 25개 생활권",
       11.5, False, INK, line=1.6)
    tb(sl, M, 4.55, W - 2 * M, .3,
       "왜 이 단계가 중요한가 — 생태학적 오류(ecological fallacy)의 사전 차단", 14, True, INK)
    tb(sl, M, 4.92, W - 2 * M, .78,
       "공간 키가 다른 데이터를 억지로 맞추려면 인구 비례 안분 같은 가정이 필요한데, 그 가정이 틀리면 지표 전체가 조용히 오염된다.\n"
       "격차 분석은 '어느 동이 더 나쁜가'를 다투는 작업이라 배분 오차가 곧 결론의 오차가 된다. 정밀도를 조금 포기하는 대신 "
       "배분 가정을 아예 쓰지 않는 쪽을 택했다.",
       12, False, MUTE, line=1.5)
    grp = cbi.groupby("ztype")["zone"].apply(lambda s: ", ".join(s.head(6))).reset_index()
    grp.columns = ["권역유형", "소속 생활권"]
    tablette(sl, grp, M, 5.62, W - 2 * M, colw=[1.2, 6.0], size=9.5, maxrows=5, rh=.30)
    footer(sl, n)

    n = _part2(prs, S, cbi, ew, sites, prov, n)
    return prs, S, cbi, ew, sites, prov, n


def _part2(prs, S, cbi, ew, sites, prov, n):
    mode_xsec = S.get("모델_모드") == "cross-section"

    # ══ 8. 진단 방법론 ═══════════════════════════════════════
    n += 1
    sl = slide(prs, "METHOD · 02", "진단 엔진 — 엔트로피 가중 균형발전지수(CBI)",
               "가중치를 연구자가 정하지 않는다. 데이터의 변별력에서 가중치를 유도한다.")
    doms = [("인구활력", "5년 인구증감률 · 청년비율 · 고령비율", ACC),
            ("경제활력", "인구천명당 사업체수 · 상권 순증감 · 폐업률 · 업종 다양성", AMBER),
            ("생활SOC", "생활SOC 접근성(중력모형) · 인구천명당 SOC수", BLUE),
            ("주거", "빈집 추정률 · 노후(30년+) 건물 비율", GREEN),
            ("이동성", "km²당 버스정류장 수", C(0x8A, 0x7F, 0xB5))]
    tb(sl, M, 2.5, 6.1, .3, f"5대 도메인 · {S['지표수']}개 지표", 14, True, INK)
    for i, (d, ind, col) in enumerate(doms):
        y = 2.92 + i * .62
        rect(sl, M, y, .075, .48, col, radius=False)
        tb(sl, M + .24, y + .02, 1.3, .28, d, 12.5, True, INK)
        tb(sl, M + 1.66, y + .04, 4.35, .42, ind, 10.5, False, MUTE, line=1.35)
    rect(sl, M + 6.35, 2.5, 5.55, 3.55, LIGHT)
    tb(sl, M + 6.65, 2.74, 5, .3, "왜 엔트로피 가중법인가", 14, True, INK)
    tb(sl, M + 6.65, 3.14, 5, .6,
       "pⱼ = zⱼ / Σzⱼ    eⱼ = −k·Σ pⱼ ln pⱼ    wⱼ ∝ (1 − eⱼ)",
       12.5, True, ACC, line=1.4, font="Consolas")
    tb(sl, M + 6.65, 3.72, 5, 2.1,
       "생활권 간 차이를 잘 벌리는 지표일수록(=엔트로피가 낮을수록) 큰 가중치를 받는다.\n\n"
       "· 전문가 설문·AHP 방식과 달리 사람이 개입하지 않아 누가 돌려도 같은 값이 나온다\n"
       "· 지표가 추가·제외돼도 가중치가 자동 재배분되어 지수가 계속 유효하다\n"
       "· 심사·감사 과정에서 '왜 이 가중치인가'를 수식 한 줄로 답할 수 있다",
       11.5, False, MUTE, line=1.55)
    rect(sl, M, 6.12, W - 2 * M, .78, C(0xF0, 0xF7, 0xF2))
    tb(sl, M + .3, 6.32, 11.5, .45,
       "정직성 장치 — 데이터가 없어 산출 불가한 지표는 임의값으로 채우지 않고 지수에서 제외하며,\n"
       "제외 사실을 로그와 산출물에 남긴다. 결측을 평균으로 메워 순위를 왜곡하는 흔한 실수를 구조적으로 막는다.",
       11.5, False, INK, line=1.5)
    footer(sl, n)

    # ══ 9. 진단 결과 ═════════════════════════════════════════
    n += 1
    sl = slide(prs, "FINDING · 01", "진단 결과 — 25개 생활권 CBI",
               f"최상위 {S['최고']}점, 최하위 {S['최저']}점. 같은 시(市) 안의 격차다.")
    fit(sl, "01_CBI_랭킹.png", M, 2.28, 7.0, 4.5)
    x = M + 7.3
    kpi_card(sl, x, 2.32, 4.55, 1.34, "신도심 ↔ 원도심 격차",
             f'{S["격차배율"]}배', f'{S["신도심_CBI"]} vs {S["원도심_CBI"]}', ACC)
    kpi_card(sl, x, 3.76, 4.55, 1.34, "CBI 지니계수",
             f'{S["지니계수"]}', f'변동계수 {S["변동계수"]}%', AMBER)
    kpi_card(sl, x, 5.20, 4.55, 1.34, "군집 유형화",
             f'{S["군집수"]}개 유형', f'실루엣 계수 {S["실루엣"]}', BLUE)
    tb(sl, x, 6.62, 4.55, .32,
       "쇠퇴단계는 CBI 순위가 아니라 다차원 패턴으로 분류된다.",
       10, False, MUTE, line=1.35)
    footer(sl, n)

    # ══ 10. 핵심 인사이트 ════════════════════════════════════
    n += 1
    sl = slide(prs, "FINDING · 02", "같은 '어려움'이라도 필요한 도움이 다릅니다",
               "원도심과 농촌면은 지수가 모두 낮게 나왔습니다. 그런데 항목을 나눠 보니 사정이 서로 달랐습니다.")
    fit(sl, "02_도메인_히트맵.png", M, 2.36, 7.35, 3.5)
    x = M + 7.65
    rect(sl, x, 2.36, 4.2, 1.72, C(0xFD, 0xF2, 0xF2))
    tb(sl, x + .26, 2.58, 3.7, .28, "원도심", 13.5, True, ACC)
    tb(sl, x + .26, 2.92, 3.7, 1.05,
       "시설은 이미 갖춰져 있습니다.\n어려운 쪽은 상권과 인구였습니다.\n"
       "→ 새로 짓기보다, 비어 있는 공간을 다시 쓰는 방향이 맞지 않을까 합니다",
       11, False, INK, line=1.5)
    rect(sl, x, 4.22, 4.2, 1.72, C(0xF3, 0xF0, 0xF9))
    tb(sl, x + .26, 4.44, 3.7, .28, "농촌면 · 읍지역", 13.5, True, C(0x6B, 0x5F, 0x9E))
    tb(sl, x + .26, 4.78, 3.7, 1.05,
       "가까운 시설 자체가 부족합니다.\n대중교통 여건도 넉넉하지 않습니다.\n"
       "→ 거점 복합시설과 이동 지원을 함께 보시면 좋겠습니다",
       11, False, INK, line=1.5)
    rect(sl, M, 6.08, W - 2 * M, .82, LIGHT)
    tb(sl, M + .3, 6.28, 11.5, .5,
       "여러 동네를 한 묶음으로 보면 놓치기 쉬운 부분입니다. 유형을 나눠 보면 같은 예산으로도 각 동네에 조금 더 맞는\n"
       "도움을 드릴 수 있지 않을까 합니다. 어떤 사업이 맞을지는 결국 현장을 아시는 분들의 판단이 필요한 부분입니다.",
       12, False, INK, line=1.5)
    footer(sl, n)

    # ══ 11. 격차 추이 ════════════════════════════════════════
    n += 1
    sl = slide(prs, "FINDING · 03", "격차는 수렴하지 않고 벌어지고 있다",
               "권역유형별 인구·사업체 궤적이 시간이 갈수록 분기한다. 개입하지 않으면 자동으로 좁혀지지 않는다.")
    fit(sl, "03_격차추이.png", M, 2.5, W - 2 * M, 3.5)
    rect(sl, M, 6.18, W - 2 * M, .78, C(0xFD, 0xF2, 0xF2))
    tb(sl, M + .3, 6.38, 11.5, .45,
       "이 그림이 조기경보 엔진의 출발점이다 — 추세가 이미 갈라졌다면, 다음에 갈라질 곳도 지금의 구조에서 읽어낼 수 있다.",
       12.5, True, INK, line=1.5)
    footer(sl, n)

    # ══ 12. 예측 설계 ════════════════════════════════════════
    n += 1
    sl = slide(prs, "MODEL", "미리 살펴보기 — 쇠퇴 조기경보 모델",
               "이미 어려워진 곳뿐 아니라, 앞으로 살펴보면 좋을 곳도 함께 짚어보려 만들었습니다.")
    left = [("무엇을 살펴보나요",
             "t년 정보만으로 t+3년 상황을 미리 봅니다. 예산 편성에 1년, 사업 착수에 1~2년이\n"
             "걸리는 점을 감안해 3년으로 잡았습니다."),
            ("모델", "LightGBM 회귀입니다. 표본이 많지 않아 규제를 강하게 걸고 얕은 트리를 썼습니다."),
            ("검증 — Leave-One-Zone-Out CV",
             "같은 생활권이 학습과 검증에 동시에 들어가지 않게 했습니다. 가까운 동네끼리는\n"
             "닮아 있어 무작위로 나누면 성능이 실제보다 좋아 보이기 때문입니다."),
            ("이유까지 함께 — SHAP",
             "동네마다 어떤 항목 때문에 그렇게 나왔는지 나눠서 보여드립니다.")]
    bullets(sl, M, 2.5, 6.3, left, size=12, gap=.38)
    rect(sl, M + 6.75, 2.55, 5.15, 3.6, WHITE, LINE)
    tb(sl, M + 7.05, 2.8, 4.5, .3, "성능 (현재 데이터 기준)", 13.5, True, INK)
    _ok = S.get("예측엔진", True)
    _hit = S.get("방향적중률")
    perf = [("검증 방식", "Leave-One-Zone-Out CV"),
            ("예측 목표", "3년 후 위험도 변화량(Δ)"),
            ("방향 적중률", f'{_hit:.0f}%' if (_ok and _hit == _hit) else "— (데이터 대기)"),
            ("MAE", f'{S["모델_MAE"]} (베이스라인 {S["베이스라인_MAE"]})' if _ok else "—"),
            ("베이스라인 대비", f'{S["MAE_개선율"]:+.1f}%' if _ok else "—"),
            ("R²", f'{S["모델_R2"]}' if _ok else "—")]
    for i, (k, v) in enumerate(perf):
        y = 3.16 + i * .42
        tb(sl, M + 7.05, y, 2.0, .3, k, 11.5, False, MUTE)
        tb(sl, M + 9.15, y, 2.5, .3, v, 12, True, INK)
    tb(sl, M + 7.05, 5.68, 4.6, .5,
       "베이스라인은 '지금과 달라지지 않는다(Δ=0)'는 예측입니다.\n"
       "지금은 방향을 주로 맞히는 수준이며, 변화의 크기까지는 아직입니다(R²가 낮은 이유).",
       10, False, MUTE, line=1.4)
    rect(sl, M, 6.32, W - 2 * M, .62, LIGHT)
    tb(sl, M + .3, 6.48, 11.5, .35,
       "데이터 상황에 따라 모드가 자동 전환된다 — 개·폐업 시계열이 확보되면 패널 예측(t→t+3)으로, "
       "없으면 구조지표 횡단면 모드로 돌아간다.", 11.5, False, INK, line=1.45)
    footer(sl, n)

    # ══ 13. 예측 결과 ════════════════════════════════════════
    n += 1
    sl = slide(prs, "FINDING · 04", "예측 결과 — 다음 3년, 어디를 봐야 하는가",
               "대각선 위쪽은 현재보다 위험이 더 올라갈 것으로 예측된 생활권이다.")
    fit(sl, "04_조기경보.png", M, 2.3, 5.6, 4.35)
    fit(sl, "05_SHAP_요인분해.png", M + 5.8, 2.3, 6.1, 3.1)
    if not len(ew):
        rect(sl, M + 5.8, 5.5, 6.1, .9, C(0xFD, 0xF8, 0xEC))
        tb(sl, M + 6.05, 5.68, 5.6, .55,
           "예측 엔진 미실행 — " + S.get("예측미실행사유", ""), 10.5, False, INK, line=1.4)
        footer(sl, n); return _part3(prs, S, cbi, ew, sites, prov, n)
    e = ew.head(5).copy()
    e["risk"] = e["risk"].round(0).astype(int); e["risk_pred"] = e["risk_pred"].round(0).astype(int)
    e["risk_delta"] = e["risk_delta"].round(0).astype(int).map(lambda v: f"{v:+d}")
    e.columns = ["생활권", "현재", "3년 후", "변화"]
    tb(sl, M + 5.8, 5.42, 6.1, .26, "3년 뒤 위험 상위 5개 생활권", 11.5, True, INK)
    tablette(sl, e, M + 5.8, 5.74, 6.1, colw=[2.2, 1.2, 1.2, 1.2], size=9.5, maxrows=5,
             rh=.28)
    footer(sl, n)
    return _part3(prs, S, cbi, ew, sites, prov, n)


def _part3(prs, S, cbi, ew, sites, prov, n):

    # ══ 14. 처방 설계 ════════════════════════════════════════
    n += 1
    sl = slide(prs, "MODEL", "순서 정하기 — 생활SOC 우선순위 계산",
               "예산을 얼마나 쓸지가 아니라, 한 곳을 놓는다면 어디가 가장 많은 분께 닿을지를 계산했습니다.")
    tb(sl, M, 2.5, 6.2, .3, "MCLP — 최대커버링 입지문제", 14, True, INK)
    tb(sl, M, 2.88, 6.2, .5,
       "max  Σ  dᵢ · yᵢ      s.t.  yᵢ ≤ Σ xⱼ ,  Σ xⱼ = p\n"
       "        i∈수요격자                    j∈Nᵢ",
       12, True, ACC, line=1.5, font="Consolas")
    bullets(sl, M, 3.52, 6.2, [
        ("수요 dᵢ — 취약가중 인구",
         "격자 인구에 고령화율과 CBI 열위를 가중한다. 같은 1명이라도 취약지역의 1명을 더 크게 센다."),
        ("커버 Nᵢ — 반경 1.0km",
         "도보 15분권. '15분 도시' 개념과 맞추고 행정경계를 넘는 이용도 반영한다."),
        ("후보지 xⱼ — 빈집·노후 밀집 격자",
         "새 땅을 사는 대신 이미 비어 있는 공간을 후보로 둔다. 도시재생 사업과 바로 접속된다."),
        ("탐욕 근사의 보장",
         "submodular 목적함수이므로 탐욕해가 최적해의 (1−1/e)≈63% 이상을 보장한다."),
    ], size=11.5, gap=.38)
    rect(sl, M + 6.65, 2.5, 5.25, 3.05, C(0xF0, 0xF7, 0xF2))
    tb(sl, M + 6.95, 2.74, 4.7, .3, "형평성 제약을 넣은 이유", 13.5, True, GREEN)
    tb(sl, M + 6.95, 3.14, 4.7, 2.2,
       "계산만 그대로 두면 효율이 가장 높은 한 동네에 전부 몰립니다. 수식으로는 맞지만 "
       "균형발전이라는 취지에는 맞지 않는다고 생각했습니다.\n\n"
       "그래서 생활권당 최대 2개소로 제한했습니다. 총 효과는 조금 줄지만 "
       "여러 동네에 나누어 닿습니다.\n\n"
       "이 값은 조정 가능하니, 시 사정에 맞게 바꿔 쓰시면 됩니다.",
       11, False, INK, line=1.55)
    rect(sl, M + 6.65, 5.68, 5.25, 1.22, LIGHT)
    tb(sl, M + 6.95, 5.88, 4.7, .85,
       "결과는 '순위 · 생활권 · 시설유형 · 새로 닿는 인구 · 개선폭'이 적힌 표 한 장입니다.\n"
       "검토 자료로 참고하실 수 있으면 좋겠습니다.", 11, False, INK, line=1.5)
    footer(sl, n)

    # ══ 15. 처방 결과 ════════════════════════════════════════
    n += 1
    sl = slide(prs, "FINDING · 05", "처방 결과 — 투자 우선순위",
               f"상위 {S['추천입지수']}개소 우선 투자 시 취약수요 커버리지 +{S['커버리지개선']}%p, "
               f"신규 수혜인구 약 {S['신규수혜인구']:,}명.")
    fit(sl, "07_투자입지_지도.png", M, 2.42, 5.3, 4.3)
    st = sites[["순위", "생활권", "시설유형", "신규수혜인구", "커버리지개선률"]].head(10).copy()
    st["신규수혜인구"] = st["신규수혜인구"].map("{:,}".format)
    st["커버리지개선률"] = st["커버리지개선률"].map(lambda v: f"+{v}%p")
    st.columns = ["순위", "생활권", "시설유형", "신규 수혜인구", "커버리지 개선"]
    tb(sl, M + 5.6, 2.42, 6.3, .3, "AI 추천 투자 순위 (상위 10)", 13.5, True, INK)
    tablette(sl, st, M + 5.6, 2.8, 6.3, colw=[.7, 1.4, 1.2, 1.6, 1.4], size=10, maxrows=10)
    footer(sl, n)

    # ══ 16. 예산 배분 곡선 ═══════════════════════════════════
    n += 1
    sl = slide(prs, "FINDING · 06", "어디서 멈춰도 근거가 남는다",
               "한계효용 체감 곡선이 있으면 예산이 깎여도 '몇 개소까지가 합리적인가'를 답할 수 있다.")
    fit(sl, "08_커버리지_곡선.png", M, 2.42, 7.2, 3.5)
    fit(sl, "06_SOC_사각지대.png", M + 7.5, 2.42, 4.4, 3.5)
    rect(sl, M, 6.1, W - 2 * M, .82, LIGHT)
    tb(sl, M + .3, 6.3, 11.5, .5,
       "실무적 가치 — 예산 심의에서 가장 자주 나오는 질문은 '이걸 왜 이 순서로 하느냐'와 '반만 하면 어떻게 되느냐'다.\n"
       "이 곡선은 두 질문에 동시에 답한다. 좌측 그림의 좌상단 영역(고령↑·SOC↓)이 최우선 사각지대다.",
       12, False, INK, line=1.5)
    footer(sl, n)

    # ══ 신설. 이 제안이 닿았으면 하는 분들 ═══════════════════
    n += 1
    sl = slide(prs, "FOR WHOM", "결국, 이런 분들께 닿았으면 합니다",
               "지수와 모델은 수단일 뿐입니다. 저희가 계속 떠올린 것은 아래 네 분의 하루였습니다.")
    lo3 = cbi.nsmallest(3, "CBI").index.tolist()
    old_top = cbi.nlargest(1, "aging_ratio")
    people = [
        ("원도심에 오래 사신 어르신", ACC,
         f"{lo3[0] if lo3 else '원도심'} 일대는 젊은 이웃이 하나둘 떠나고 상가도 비었습니다.\n"
         "가게와 병원이 멀어질수록, 걸어서 하실 수 있는 일이 줄어듭니다."),
        ("농촌면에서 아이 키우시는 부모님", C(0x8A, 0x7F, 0xB5),
         f"고령 비율이 가장 높은 곳은 {old_top.index[0]}으로 "
         f"{old_top['aging_ratio'].iloc[0]:.0f}%였습니다.\n"
         "어린이집·도서관이 가까이 없으면 그 부담은 온전히 가정의 몫이 됩니다."),
        ("원도심에서 가게를 지켜오신 상인", AMBER,
         "오가는 분이 줄면 매출이 줄고, 문 닫는 가게가 늘면 오가는 분이 더 줍니다.\n"
         "이 고리를 일찍 알아차릴 수 있다면, 도움드릴 방법도 더 많아집니다."),
        ("신도심에서 자리 잡은 이웃", BLUE,
         "지금 여건이 좋은 곳도 언젠가는 나이가 듭니다.\n"
         "미리 살펴보는 습관은 결국 천안시민 모두를 위한 준비라고 생각합니다."),
    ]
    for i, (who, col, body) in enumerate(people):
        x = M + (i % 2) * 6.1
        y = 2.5 + (i // 2) * 1.72
        rect(sl, x, y, 5.75, 1.5, WHITE, LINE)
        rect(sl, x, y, .075, 1.5, col, radius=False)
        tb(sl, x + .28, y + .2, 5.1, .3, who, 13.5, True, col)
        tb(sl, x + .28, y + .6, 5.1, .8, body, 11, False, MUTE, line=1.5)
    rect(sl, M, 6.06, W - 2 * M, .86, C(0xF0, 0xF7, 0xF2))
    tb(sl, M + .3, 6.26, 11.5, .5,
       "저희는 천안시 사정을 현장에서 겪어보지 못한 참가자입니다. 다만 공개된 데이터로 할 수 있는 만큼은 성실하게 살펴보았습니다.\n"
       "부족한 부분은 현장을 아시는 분들이 채워주시길 바라며, 이 자료가 그 논의의 출발점이 될 수 있다면 더 바랄 것이 없겠습니다.",
       12, False, INK, line=1.5)
    footer(sl, n)

    # ══ 17. 정책 활용 시나리오 ═══════════════════════════════
    n += 1
    sl = slide(prs, "APPLICATION", "실무에 이렇게 보탬이 되면 좋겠습니다",
               "따로 시간을 내어 다뤄야 하는 연구물이 아니라, 하시던 업무 곁에 놓고 참고하실 수 있는 자료로 준비했습니다.")
    users = [("도시재생 담당", "사업지 검토 · 공모 준비",
              "지수가 낮은 곳과 3년 뒤 살펴볼 곳을 겹쳐 후보를 좁히실 때,\n"
              "공모의 '쇠퇴도 진단' 항목에 근거 수치로."),
             ("복지·보건 담당", "생활시설 위치 검토",
              "우선순위표의 좌표를 검토 후보 중 하나로,\n"
              "왜 그 동네인지 설명하실 때 항목별 근거로."),
             ("예산 담당", "배분 검토 자료",
              "몇 개소까지면 어느 정도 효과인지 곡선으로 비교하실 때,\n"
              "예산 조정이 필요할 때 판단의 참고로."),
             ("기획·평가 담당", "시행 이후 확인",
              "이듬해 같은 스크립트로 지수 변화를 이어 보실 때,\n"
              "하신 사업이 수치로도 나아졌는지 확인하실 때.")]
    for i, (who, what, how) in enumerate(users):
        x = M + (i % 2) * 6.1
        y = 2.5 + (i // 2) * 1.72
        rect(sl, x, y, 5.75, 1.5, WHITE, LINE)
        rect(sl, x, y, .075, 1.5, [ACC, BLUE, GREEN, AMBER][i], radius=False)
        tb(sl, x + .28, y + .18, 2.4, .28, who, 12.5, True, [ACC, BLUE, GREEN, AMBER][i])
        tb(sl, x + .28, y + .5, 5.1, .28, what, 13.5, True, INK)
        tb(sl, x + .28, y + .84, 5.1, .6, how, 10.8, False, MUTE, line=1.45)
    rect(sl, M, 6.06, W - 2 * M, .86, C(0xF0, 0xF7, 0xF2))
    tb(sl, M + .3, 6.24, 11.5, .55,
       "쓰시는 방법 — 연 1회 공공데이터를 새로 받아 스크립트를 실행하시면 모든 자료가 다시 만들어집니다.\n"
       "별도 시스템이나 데이터 구매 없이, 담당자 한 분이 반나절이면 그 해 자료를 준비하실 수 있습니다.",
       12, False, INK, line=1.5)
    footer(sl, n)

    # ══ 18. 기대효과 & 로드맵 ════════════════════════════════
    n += 1
    sl = slide(prs, "IMPACT", "기대효과와 실행 로드맵", None)
    tb(sl, M, 2.28, 6.0, .3, "기대효과", 14.5, True, INK)
    eff = [("시민께 닿는 범위", "효과가 큰 곳부터 살펴보면, 같은 예산으로도 더 많은 분께 닿을 수 있습니다"),
           ("살펴보는 시점", "어려워진 뒤가 아니라, 3년 앞을 미리 짚어볼 수 있습니다"),
           ("설명하실 때", "판단 근거를 항목별 수치로 정리해 드려 설명 부담을 조금 덜어드립니다"),
           ("드는 비용", "데이터 구매나 시스템 없이, 스크립트 재실행만으로 매년 갱신됩니다")]
    for i, (k, v) in enumerate(eff):
        y = 2.7 + i * .82
        rect(sl, M, y, .075, .62, ACC, radius=False)
        tb(sl, M + .24, y, 5.6, .28, k, 12.5, True, INK)
        tb(sl, M + .24, y + .3, 5.6, .4, v, 11, False, MUTE, line=1.4)
    tb(sl, M + 6.4, 2.28, 5.5, .3, "실행 로드맵", 14.5, True, INK)
    road = [("1단계 · 즉시", "본 스크립트로 천안시 25개 생활권 기준선 진단 확정", GREEN),
            ("2단계 · 3개월", "천안시 내부 행정데이터(인허가 원장·건축물대장·복지수급) 연계로 지표 정밀화", BLUE),
            ("3단계 · 6개월", "500m 격자 단위로 공간해상도 상향 — 표본 확대로 예측모델 성능 개선", AMBER),
            ("4단계 · 1년", "천안시 행정포털 대시보드 내재화, 연 1회 자동 갱신 체계 정착", ACC)]
    for i, (t_, d_, col) in enumerate(road):
        y = 2.7 + i * .82
        rect(sl, M + 6.4, y, .075, .62, col, radius=False)
        tb(sl, M + 6.64, y, 5.2, .28, t_, 12.5, True, col)
        tb(sl, M + 6.64, y + .3, 5.2, .4, d_, 11, False, MUTE, line=1.4)
    rect(sl, M, 6.02, W - 2 * M, .9, LIGHT)
    tb(sl, M + .3, 6.22, 11.5, .55,
       "다른 지역에도 — 공간단위 정의(config.py)만 바꾸면 아산·청주 등에도 같은 방식으로 쓸 수 있습니다.\n"
       "원도심 공동화는 천안만의 일이 아니어서, 도움이 된다면 충남 시군이 함께 쓰는 자료가 되어도 좋겠습니다.",
       12, False, INK, line=1.5)
    footer(sl, n)

    # ══ 19. 한계와 보완 ══════════════════════════════════════
    n += 1
    sl = slide(prs, "LIMITATION", "저희가 아직 못 한 부분을 먼저 말씀드립니다",
               "행정에 쓰이는 자료인 만큼, 어디까지 믿고 쓰실 수 있는지 분명히 알려드리는 것이 도리라고 생각했습니다.")
    lim = [("예측 성능", "3년 뒤 악화·개선 '방향'은 어느 정도 맞히지만, 변화의 '크기'까지는 아직 부족합니다.",
            "위험도 수준은 자기상관이 강해 '그대로 간다'는 예측이 이미 강력한 기준선입니다. "
            "현재는 인구 축만 실데이터라, 상권 개·폐업 시계열이 더해지면 나아질 것으로 봅니다."),
           ("표본 크기", f"분석 단위가 {S['생활권수']}개 생활권이라 머신러닝 표본으로는 작습니다.",
            "500m 격자 단위로 내리면 표본이 수천 개로 늘어납니다. 지금은 규제를 강하게 걸고 "
            "Leave-One-Zone-Out CV와 베이스라인 비교로 과적합 여부를 확인하고 있습니다."),
           ("공간 해상도", "생활권 안쪽의 차이(같은 동 안의 편차)는 잡히지 않습니다.",
            "집계구·격자 데이터로 바꾸면 해결됩니다. 공간단위 정의만 교체하면 되도록 만들어 두었습니다."),
           ("인과 아닌 상관", "이 모델은 변화를 미리 살펴볼 뿐, 사업의 효과를 증명하지는 못합니다.",
            "시행 전후 이중차분(DID) 설계를 덧붙이면 효과 검증까지 넓힐 수 있습니다.")]
    for i, (k, prob, fix) in enumerate(lim):
        y = 2.5 + i * 1.12
        rect(sl, M, y, W - 2 * M, .98, WHITE, LINE)
        rect(sl, M, y, .075, .98, AMBER, radius=False)
        tb(sl, M + .3, y + .14, 2.3, .28, k, 12.5, True, AMBER)
        tb(sl, M + 2.75, y + .14, 8.8, .28, prob, 12, True, INK)
        tb(sl, M + 2.75, y + .46, 8.8, .42, "→ " + fix, 10.8, False, MUTE, line=1.4)
    footer(sl, n)

    # ══ 20. 재현성·윤리 ══════════════════════════════════════
    n += 1
    sl = slide(prs, "INTEGRITY", "재현성과 데이터 윤리", None)
    boxes = [("재현성", GREEN, [
        "전 과정이 공개 스크립트로 실행 — `python3 src/run_all.py` 한 줄",
        "난수 시드 고정(20260831) — 군집·모델 결과가 실행마다 동일",
        "그림 8종·분석표 9종·대시보드가 같은 실행에서 함께 생성됨",
        "지표 정의·가중치·모델 파라미터가 전부 코드에 명시"]),
        ("데이터 윤리", BLUE, [
        "정식 공공데이터 포털에서 내려받은 파일만 사용 — 크롤링 없음",
        "상업 목적 민간데이터 무단 사용 없음",
        "전 지표가 생활권 단위 집계값 — 개인식별정보 미취급",
        "출처·수집시점·컬럼을 `00_데이터_출처.csv` 에 자동 기록"])]
    for i, (title_, col, items) in enumerate(boxes):
        x = M + i * 6.1
        rect(sl, x, 2.32, 5.75, 2.72, WHITE, LINE)
        rect(sl, x, 2.32, 5.75, .075, col, radius=False)
        tb(sl, x + .3, 2.56, 4.5, .3, title_, 14.5, True, col)
        for j, it in enumerate(items):
            tb(sl, x + .3, 3.0 + j * .52, .2, .28, "·", 13, True, col)
            tb(sl, x + .52, 3.0 + j * .52, 5.0, .45, it, 11, False, INK, line=1.4)
    rect(sl, M, 5.26, W - 2 * M, 1.15, C(0xFD, 0xF2, 0xF2))
    tb(sl, M + .32, 5.48, 11.4, .32, "실데이터 / 예시데이터 구분 장치", 13.5, True, ACC)
    tb(sl, M + .32, 5.84, 11.4, .5,
       "파이프라인은 각 지표가 실데이터인지 예시(illustrative) 값인지를 매 실행마다 판정해 로그·대시보드·본 기획서에 표기한다.\n"
       "예시 데이터를 실측치로 오인해 제출하는 사고를 막기 위한 장치이며, 실데이터를 넣으면 해당 표기가 자동으로 사라진다.",
       11.5, False, INK, line=1.5)
    tb(sl, M, 6.62, W - 2 * M, .3,
       f"현재 상태: 실데이터 지표 {S['실데이터지표']}"
       + ("  ·  전 지표 실데이터로 제출 가능" if S["전체실데이터"]
          else "  ·  제출 전 실데이터 투입 후 재실행 필요"),
       12, True, GREEN if S["전체실데이터"] else ACC)
    footer(sl, n)
    return n


def save(team="○○팀", members="홍길동"):
    prs, *_ = build(team, members)
    out = DELIV / "기획서_천안균형발전나침반_CBC.pptx"
    prs.save(str(out))
    print(f"    · {out.relative_to(DELIV.parent)}  ({out.stat().st_size / 1024:.0f} KB, "
          f"{len(prs.slides.__iter__.__self__._sldIdLst)}장)")
    return out
